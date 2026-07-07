import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml import parse_xml
import urllib.request

def add_fade_transition(slide):
    # Access the XML of the slide
    slide_xml = slide.element
    
    # Check if a transition already exists
    transition_el = slide_xml.xpath('./p:transition')
    if transition_el:
        slide_xml.remove(transition_el[0])
        
    # Add fade transition
    fade_transition_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>'
    fade_el = parse_xml(fade_transition_xml)
    
    # Insert transition element before timing or extLst if they exist, or at the end
    timing_els = slide_xml.xpath('./p:timing')
    if timing_els:
        timing_els[0].addprevious(fade_el)
    else:
        slide_xml.append(fade_el)

def apply_theme(slide):
    # Set background color to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Set all text to appropriate colors and fonts
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        
        # If it's a title (usually the first shape or highest)
        is_title = shape == slide.shapes[0]
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                if is_title:
                    run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
                    run.font.size = Pt(30)
                else:
                    run.font.color.rgb = RGBColor(50, 50, 50) # Dark Gray
                    run.font.size = Pt(18)

def add_footer(slide, slide_num):
    # Add footer
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "MainCrafts Technology | AI & ML Internship"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.font.name = 'Arial'
    
    # Add slide number
    txBox2 = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(0.5), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(slide_num)
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(128, 128, 128)
    p2.font.name = 'Arial'

def create_table(slide, rows, cols, left, top, width, height, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            # Style
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Arial'
                    if r == 0:
                        run.font.size = Pt(20)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.bold = True
                    else:
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Background color for header and rows
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = RGBColor(0, 51, 102) # Blue
            elif r % 2 == 1:
                fill.fore_color.rgb = RGBColor(240, 248, 255) # Light Blue
            else:
                fill.fore_color.rgb = RGBColor(255, 255, 255) # White
                
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = width

def process_presentation():
    pptx_path = 'presentation/Task1_Presentation.pptx'
    prs = Presentation(pptx_path)
    
    try:
        with open('plots/metrics.json', 'r') as f:
            metrics = json.load(f)
    except:
        metrics = {"MAE": 0.0, "MSE": 0.0, "RMSE": 0.0, "R2 Score": 0.0}
        
    for i, slide in enumerate(prs.slides):
        add_fade_transition(slide)
        apply_theme(slide)
        add_footer(slide, i + 1)
        
        # Determine slide based on title
        title = slide.shapes.title.text if slide.shapes.title else ""
        
        if "House Price Prediction" in title:
            # 1. Title Slide
            # Remove old subtitle and add new one
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape != slide.shapes.title:
                    sp = shape.element
                    sp.getparent().remove(sp)
            
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            text = [
                "Prepared by: Shankramma Baganal",
                "Department of Information Science & Engineering",
                "SDM College of Engineering and Technology, Dharwad",
                "AI & ML Internship – MainCrafts Technology",
                "Submission Date: July 2026"
            ]
            for j, line in enumerate(text):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER
                
        elif "Dataset Overview" in title:
            # 4. Dataset Overview -> Table
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape != slide.shapes.title:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
            data = [
                ["Property", "Details"],
                ["Dataset Name", "California Housing Dataset"],
                ["Source", "Scikit-Learn"],
                ["Number of Samples", "20,640"],
                ["Number of Features", "8"],
                ["Target Variable", "Median House Value"]
            ]
            create_table(slide, 6, 2, Inches(1), Inches(2), Inches(8), Inches(3.5), data, [Inches(3), Inches(5)])
            
        elif "Machine Learning Workflow" in title:
            # 5. Workflow Slide -> SmartArt / Flow diagram
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape != slide.shapes.title:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
            steps = [
                "Dataset", "Data Preprocessing", "EDA", 
                "Train-Test Split", "Linear Regression", 
                "Prediction", "Evaluation"
            ]
            
            # Draw flow chart
            y_start = Inches(1.5)
            x_start = Inches(3.5)
            box_width = Inches(3)
            box_height = Inches(0.5)
            spacing = Inches(0.7)
            
            for j, step in enumerate(steps):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    x_start, y_start + j*spacing, 
                    box_width, box_height
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 51, 102)
                line = shape.line
                line.color.rgb = RGBColor(255, 255, 255)
                
                tf = shape.text_frame
                tf.text = step
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.name = 'Arial'
                
                # Draw arrow
                if j < len(steps) - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        x_start + box_width/2 - Inches(0.1),
                        y_start + j*spacing + box_height,
                        Inches(0.2), Inches(0.2)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(0, 102, 204)
                    arrow.line.color.rgb = RGBColor(0, 102, 204)
                    
        elif "Model Evaluation Metrics" in title:
            # 6. Evaluation Metrics -> Table
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape != slide.shapes.title:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
            data = [
                ["Metric", "Value"],
                ["Mean Absolute Error (MAE)", f"{metrics.get('MAE', 0):.4f}"],
                ["Mean Squared Error (MSE)", f"{metrics.get('MSE', 0):.4f}"],
                ["Root Mean Sq. Error (RMSE)", f"{metrics.get('RMSE', 0):.4f}"],
                ["R-squared (R2) Score", f"{metrics.get('R2 Score', 0):.4f}"]
            ]
            create_table(slide, 5, 2, Inches(1.5), Inches(2.5), Inches(7), Inches(2.5), data, [Inches(4.5), Inches(2.5)])
            
        elif "Results" in title:
            # 7. Results Slide -> Short observation below graph
            txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Observation: The model predictions align closely with actual prices along the diagonal, though some variance exists for extreme values."
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.font.name = 'Arial'
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            
        elif "Conclusion & Future Scope" in title:
            # 8. Conclusion Slide -> Key Learning Outcomes
            slide.shapes.title.text = "Conclusion & Key Learning Outcomes"
            
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape != slide.shapes.title:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            outcomes = [
                "1. Successfully developed an end-to-end Machine Learning pipeline.",
                "2. Gained practical experience in Exploratory Data Analysis (EDA) and feature correlation.",
                "3. Implemented and evaluated a Multiple Linear Regression model using Scikit-Learn.",
                "4. Understood the impact of regression evaluation metrics (MAE, MSE, RMSE, R²).",
                "5. Improved skills in visualizing data insights and model predictions."
            ]
            
            for j, line in enumerate(outcomes):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.font.name = 'Arial'
                # p.level = 0 # for bullets, but using numbered text instead
                p.space_after = Pt(10)

    # Note: Downloading simple placeholders for professional icons
    try:
        # We will add simple text shapes as icons if download fails, but let's try just downloading a python logo
        # For an internship presentation, placing small styled circles as 'icons' around the title slide can work,
        # but let's just stick to clean formatting if icons aren't easily downloaded.
        # Let's add a decorative blue bar at the top of every slide
        for slide in prs.slides:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
            shape.line.color.rgb = RGBColor(0, 51, 102)
    except Exception as e:
        print("Error adding decor:", e)
        
    prs.save(pptx_path)
    print("Presentation successfully updated!")

if __name__ == '__main__':
    process_presentation()
